"""Text and table extraction from documents (PDF, DOCX, PPTX, HTML).

Uses lightweight libraries:
- pdfplumber for PDF
- python-docx for DOCX
- python-pptx for PPTX
- beautifulsoup4 for HTML
"""


import io
from typing import Any, Optional

from pydantic import BaseModel, Field


class ParsedDocument(BaseModel):
    """Result of parsing an unstructured document."""

    tables: list[list[dict[str, str]]] = Field(default_factory=list)
    text_blocks: list[str] = Field(default_factory=list)
    metadata: dict[str, Any] = Field(default_factory=dict)


def parse_document(file_bytes: bytes, filename: str) -> ParsedDocument:
    """Parse a document file, extracting tables and text blocks."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    parsers = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
        "html": _parse_html,
        "htm": _parse_html,
    }

    parser = parsers.get(ext)
    if not parser:
        return ParsedDocument(text_blocks=[file_bytes.decode("utf-8", errors="replace")])

    return parser(file_bytes, filename)


# ── PDF ─────────────────────────────────────────────────────────────────────


def _parse_pdf(file_bytes: bytes, filename: str) -> ParsedDocument:
    import pdfplumber

    tables: list[list[dict[str, str]]] = []
    text_blocks: list[str] = []

    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            # Extract tables
            for table in page.extract_tables():
                rows = _raw_table_to_dicts(table)
                if rows:
                    tables.append(rows)

            # Extract text (excluding table areas)
            text = page.extract_text()
            if text and text.strip():
                text_blocks.append(text.strip())

    metadata: dict[str, Any] = {"page_count": len(pdf.pages) if hasattr(pdf, "pages") else 0}
    return ParsedDocument(tables=tables, text_blocks=text_blocks, metadata=metadata)


# ── DOCX ────────────────────────────────────────────────────────────────────


def _parse_docx(file_bytes: bytes, filename: str) -> ParsedDocument:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    tables: list[list[dict[str, str]]] = []
    text_blocks: list[str] = []
    title: Optional[str] = None

    # Extract paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if not title and para.style and "heading" in (para.style.name or "").lower():
            title = text
        text_blocks.append(text)

    # Extract tables
    for table in doc.tables:
        rows = _docx_table_to_dicts(table)
        if rows:
            tables.append(rows)

    metadata: dict[str, Any] = {}
    if title:
        metadata["title"] = title
    return ParsedDocument(tables=tables, text_blocks=text_blocks, metadata=metadata)


def _docx_table_to_dicts(table: Any) -> list[dict[str, str]]:
    """Convert a python-docx Table to list[dict]."""
    raw_rows = []
    for row in table.rows:
        raw_rows.append([cell.text.strip() for cell in row.cells])

    return _raw_table_to_dicts(raw_rows)


# ── PPTX ────────────────────────────────────────────────────────────────────


def _parse_pptx(file_bytes: bytes, filename: str) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    text_blocks: list[str] = []
    tables: list[list[dict[str, str]]] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                raw = []
                for row in shape.table.rows:
                    raw.append([cell.text.strip() for cell in row.cells])
                rows = _raw_table_to_dicts(raw)
                if rows:
                    tables.append(rows)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_blocks.append(text)

    return ParsedDocument(tables=tables, text_blocks=text_blocks, metadata={})


# ── HTML ────────────────────────────────────────────────────────────────────


def _parse_html(file_bytes: bytes, filename: str) -> ParsedDocument:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(file_bytes, "html.parser")
    tables: list[list[dict[str, str]]] = []
    text_blocks: list[str] = []

    # Extract tables
    for table_el in soup.find_all("table"):
        raw = []
        for tr in table_el.find_all("tr"):
            cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
            if cells:
                raw.append(cells)
        rows = _raw_table_to_dicts(raw)
        if rows:
            tables.append(rows)
        table_el.decompose()

    # Extract remaining text
    title = soup.title.string.strip() if soup.title and soup.title.string else None
    for tag in soup.find_all(["p", "h1", "h2", "h3", "h4", "li"]):
        text = tag.get_text(strip=True)
        if text:
            text_blocks.append(text)

    metadata: dict[str, Any] = {}
    if title:
        metadata["title"] = title
    return ParsedDocument(tables=tables, text_blocks=text_blocks, metadata=metadata)


# ── Shared Helper ───────────────────────────────────────────────────────────


def _raw_table_to_dicts(raw_rows: list[list[Optional[str]]]) -> list[dict[str, str]]:
    """Convert a list of row-lists (first row = headers) to list[dict]."""
    if len(raw_rows) < 2:
        return []

    headers = [str(h or f"col_{i}").strip() for i, h in enumerate(raw_rows[0])]
    # Skip empty or all-None header rows
    if not any(h for h in headers):
        return []

    rows = []
    for raw_row in raw_rows[1:]:
        if len(raw_row) != len(headers):
            continue
        row = {h: str(v or "").strip() for h, v in zip(headers, raw_row)}
        # Skip completely empty rows
        if any(v for v in row.values()):
            rows.append(row)
    return rows
